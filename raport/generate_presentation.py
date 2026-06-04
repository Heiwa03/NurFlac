#!/usr/bin/env python3
"""Generate NurFlac PowerPoint presentation for project defense."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

IMGS = '/home/heiwa/gitrepos/NurFlac/UML_Diagrams'
OUT  = '/home/heiwa/gitrepos/NurFlac/raport/NurFlac_Prezentare.pptx'

def img(name):
    return os.path.join(IMGS, name)

# ─── COLOR PALETTE ───────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy — background
C_ACCENT  = RGBColor(0xE9, 0x4C, 0x4C)   # red accent
C_LIGHT   = RGBColor(0xEA, 0xEA, 0xEA)   # near-white text
C_MUTED   = RGBColor(0xA0, 0xA8, 0xB8)   # muted grey text
C_CARD    = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy for cards
C_GREEN   = RGBColor(0x4C, 0xAF, 0x50)   # success green

# ─── PRESENTATION SETUP ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # blank layout

W = prs.slide_width
H = prs.slide_height

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=C_LIGHT,
          align=PP_ALIGN.LEFT, italic=False,
          line_spacing=None, space_before=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    if space_before is not None:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        spcBef = p._p.get_or_add_pPr().get_or_add_spcBef()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return shape

def add_run(para, text, size=16, bold=False, color=C_LIGHT, italic=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run

def rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1, x, y, w, h  # MSO_SHAPE_TYPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def accent_bar(slide, y=Inches(0.08), thickness=Inches(0.055)):
    """Horizontal red accent bar."""
    rect(slide, 0, y, W, thickness, C_ACCENT)

def add_picture_safe(slide, path, x, y, w, h=None):
    try:
        if h:
            slide.shapes.add_picture(path, x, y, w, h)
        else:
            slide.shapes.add_picture(path, x, y, w)
    except Exception as e:
        txbox(slide, f'[diagram]', x, y, w, Inches(1), size=12, color=C_MUTED)

def multiline_txbox(slide, lines, x, y, w, h, base_size=16):
    """Text box with multiple lines as (text, size, bold, color, italic) tuples."""
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, base_size, False, C_LIGHT, False)
        text, size, bold, color, italic = item
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return shape

def badge(slide, text, x, y, w=Inches(2.0), h=Inches(0.38),
          bg=C_ACCENT, fg=C_LIGHT, size=13, bold=True):
    """Small colored label badge."""
    r = rect(slide, x, y, w, h, bg)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg

def card(slide, title, body_lines, x, y, w, h,
         title_size=15, body_size=13, title_color=C_ACCENT):
    """Card with colored title and body text."""
    rect(slide, x, y, w, h, C_CARD)
    # title
    txbox(slide, title, x + Inches(0.15), y + Inches(0.1),
          w - Inches(0.3), Inches(0.4),
          size=title_size, bold=True, color=title_color)
    # body
    shape = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.52),
        w - Inches(0.3), h - Inches(0.65))
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(line, tuple):
            run.text = line[0]
            run.font.size = Pt(line[1]) if len(line) > 1 else Pt(body_size)
            run.font.bold = line[2] if len(line) > 2 else False
            run.font.color.rgb = line[3] if len(line) > 3 else C_LIGHT
        else:
            run.text = line
            run.font.size = Pt(body_size)
            run.font.color.rgb = C_LIGHT

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — INTRO
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)

# Large accent diagonal stripe (decorative)
rect(s, W - Inches(5.5), 0, Inches(5.5), H, RGBColor(0x0F, 0x15, 0x28))
rect(s, W - Inches(5.2), 0, Inches(0.06), H, C_ACCENT)

accent_bar(s)

# Bot emoji / logo area
txbox(s, '🎵', Inches(0.5), Inches(0.9), Inches(1.5), Inches(1.5), size=60)

# Project name — large
txbox(s, 'NurFlac', Inches(0.5), Inches(1.9), Inches(6.5), Inches(1.2),
      size=62, bold=True, color=C_LIGHT)

txbox(s, 'Bot Telegram pentru validarea fișierelor audio lossless',
      Inches(0.5), Inches(3.0), Inches(7.5), Inches(0.6),
      size=20, color=C_MUTED, italic=True)

# Divider line
rect(s, Inches(0.5), Inches(3.75), Inches(4.5), Inches(0.03), C_ACCENT)

# Student info
multiline_txbox(s, [
    ('Student:  ', 14, False, C_MUTED, False),
], Inches(0.5), Inches(4.0), Inches(6), Inches(0.4))

multiline_txbox(s, [
    ('Ciuc Vlada Marian', 22, True, C_LIGHT, False),
], Inches(0.5), Inches(4.3), Inches(6.5), Inches(0.55))

multiline_txbox(s, [
    ('Grupa  TI-XXX    ·    Disciplina: Tehnici și Mecanisme de Proiectare a Produselor Program', 14, False, C_MUTED, False),
], Inches(0.5), Inches(4.95), Inches(9), Inches(0.4))

# 12 patterns badge grid (right side)
badges = [
    ('Singleton', 'Factory Method'),
    ('Builder', 'Abstract Factory'),
    ('Decorator', 'Facade'),
    ('Adapter', 'Proxy'),
    ('State', 'Chain of Resp.'),
    ('Command', 'Strategy'),
]
bx = W - Inches(4.8)
by = Inches(1.0)
for row in badges:
    for ci, name in enumerate(row):
        bx_off = bx + ci * Inches(2.35)
        badge(s, name, bx_off, by, w=Inches(2.2), h=Inches(0.38),
              bg=RGBColor(0x16, 0x21, 0x3E), fg=C_MUTED, size=11, bold=False)
    by += Inches(0.50)

txbox(s, '12 Design Patterns  ·  .NET 9.0  ·  C# 13',
      W - Inches(4.8), by + Inches(0.1), Inches(4.6), Inches(0.4),
      size=12, color=C_ACCENT, bold=True)

# ═════════════════════════════════════════════════════════════════════════════
# HELPER — PATTERN SLIDE TEMPLATE
# ═════════════════════════════════════════════════════════════════════════════
def pattern_slide(num, name, category, img_file,
                  integration_text, problem_text,
                  category_color=None):
    """
    Standard layout for each of the 12 pattern slides.
    Left ~55%: diagram image
    Right ~45%: pattern name, category badge, integration idea, problem solved
    """
    if category_color is None:
        if 'Creațional' in category:
            category_color = RGBColor(0x4C, 0xAF, 0x50)
        elif 'Structural' in category:
            category_color = RGBColor(0x21, 0x96, 0xF3)
        else:
            category_color = RGBColor(0xFF, 0x98, 0x00)

    s = add_slide()
    fill_bg(s)
    accent_bar(s)

    # Slide number dot
    rect(s, Inches(0.25), Inches(0.2), Inches(0.55), Inches(0.55),
         C_ACCENT)
    txbox(s, str(num), Inches(0.25), Inches(0.2), Inches(0.55), Inches(0.55),
          size=16, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # ── LEFT — diagram (57% width) ──
    diag_w = Inches(7.2)
    diag_h = Inches(6.8)
    diag_x = Inches(0.22)
    diag_y = Inches(0.55)
    rect(s, diag_x, diag_y, diag_w, diag_h, RGBColor(0x0D, 0x14, 0x26))
    add_picture_safe(s, img(img_file), diag_x + Inches(0.08), diag_y + Inches(0.08),
                     diag_w - Inches(0.16), diag_h - Inches(0.16))

    # ── RIGHT — info panel (43% width) ──
    rx = Inches(7.65)
    ry = Inches(0.20)
    rw = W - rx - Inches(0.18)

    # Category badge
    badge(s, category, rx, ry, w=rw, h=Inches(0.38),
          bg=category_color, fg=RGBColor(0xFF, 0xFF, 0xFF), size=13, bold=True)

    # Pattern name
    txbox(s, name, rx, ry + Inches(0.48), rw, Inches(0.75),
          size=26, bold=True, color=C_LIGHT, align=PP_ALIGN.LEFT)

    # Integration idea card
    card_y = ry + Inches(1.35)
    card_h = Inches(2.5)
    card(s,
         '⚡  Integrare în sistem',
         integration_text,
         rx, card_y, rw, card_h,
         title_size=13, body_size=12, title_color=RGBColor(0x64, 0xB5, 0xF6))

    # Problem solved card
    prob_y = card_y + card_h + Inches(0.15)
    prob_h = H - prob_y - Inches(0.18)
    card(s,
         '🎯  Problema rezolvată',
         problem_text,
         rx, prob_y, rw, prob_h,
         title_size=13, body_size=12, title_color=RGBColor(0xA5, 0xD6, 0xA7))

    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDES 2–13 — 12 PATTERNS
# ═════════════════════════════════════════════════════════════════════════════

# 1. SINGLETON
pattern_slide(
    1, 'Singleton', 'Creațional',
    'Singleton_BotConfigurationManager.png',
    integration_text=[
        'BotConfigurationManager este punctul unic de',
        'acces la configurație (token, AdminIds, DB paths).',
        '',
        'Lazy<T> cu ExecutionAndPublication garantează',
        'o singură instanță chiar şi sub concurență.',
        '',
        'SpectralAnalysisEngine (motorul FFT) reutilizează',
        'acelaşi mecanism — o instanță partajată de toți',
        'analizatorii audio.',
    ],
    problem_text=[
        'Configurația trebuie citită o singură dată',
        'şi să fie accesibilă din orice componentă',
        'fără a fi pasată explicit prin tot DI-ul.',
    ]
)

# 2. FACTORY METHOD
pattern_slide(
    2, 'Factory Method', 'Creațional',
    'FactoryMethod_CommandFactory.png',
    integration_text=[
        'CommandFactory (Creator abstract) declară',
        'CreateCommand(token) — suprascris de',
        'StandardCommandFactory.',
        '',
        'Metoda non-virtuală Create() aplică automat',
        'decoratorul corect: AdminGuard sau',
        'ModerationGuard, în funcție de token.',
    ],
    problem_text=[
        'Dispatcher-ul nu trebuie să cunoască',
        'clasele concrete ale comenzilor.',
        'Adăugarea unei noi comenzi = o linie nouă',
        'în StandardCommandFactory, fără a modifica',
        'restul sistemului.',
    ]
)

# 3. BUILDER
pattern_slide(
    3, 'Builder', 'Creațional',
    'Builder_AlbumReportBuilder.png',
    integration_text=[
        'AlbumReportBuilder acumulează rezultatele',
        'validărilor fişier cu fişier (AddSuccess /',
        'AddFailure), thread-safe via C# 13 Lock.',
        '',
        'Build() produce AlbumReport cu listele',
        'acceptate/respinse.',
        '',
        'ToMarkdown() formatează raportul final',
        'pentru trimitere în Telegram.',
    ],
    problem_text=[
        'Raportul unui album are N fişiere validate',
        'asincron. Construcția sa treptată, thread-safe,',
        'nu poate fi realizată printr-un simplu constructor.',
    ]
)

# 4. ABSTRACT FACTORY
pattern_slide(
    4, 'Abstract Factory', 'Creațional',
    'AbstractFactory_AudioAnalyzerFactory.png',
    integration_text=[
        'IAudioAnalyzerFactory declară Create___Analyzer()',
        'pentru fiecare format (FLAC, WAV, ALAC, AIFF).',
        '',
        'LosslessAnalyzerFactory → analizoare pentru',
        'validare autentică lossless.',
        '',
        'LossyDetectorFactory → analizoare specializate',
        'în detectarea recodificărilor lossy.',
    ],
    problem_text=[
        'Comportamentul analizei variază cu familia de',
        'fişiere. Fără Abstract Factory, codul de validare',
        'ar fi plin de switch-uri pe format.',
    ]
)

# 5. DECORATOR
pattern_slide(
    5, 'Decorator', 'Structural',
    'Decorator_CommandGuards.png',
    integration_text=[
        'CommandFactory aplică exact un decorator:',
        '',
        '• Comenzi user → ModerationGuardDecorator',
        '  Blochează utilizatorii baniți/în timeout.',
        '  Ridică automat timeout-urile expirate.',
        '',
        '• Comenzi admin → AdminGuardDecorator',
        '  Ignoră silențios non-adminii.',
        '  Bypass-ează moderarea (admin bănuit',
        '  poate rula /ban).',
    ],
    problem_text=[
        'Gardurile de securitate trebuie injectate',
        'transparent, fără a modifica comenzile concrete.',
        'Cele două garduri sunt mutual exclusive',
        'şi rezolvă politici complet diferite.',
    ]
)

# 6. FACADE
pattern_slide(
    6, 'Facade', 'Structural',
    'Facade_AudioPipelineFacade.png',
    integration_text=[
        'AudioPipelineFacade cablează intern lanțul:',
        'Extension → MIME → Spectral',
        '',
        'Expune o singură metodă publică:',
        'ValidateAsync(AudioFileContext, ct)',
        '',
        'Apelantul (UpdateRouter, AlbumUploadState)',
        'nu cunoaşte existența handler-elor sau',
        'ordinea lor.',
    ],
    problem_text=[
        'Validarea audio implică 3 handler-e cu',
        'dependențe diferite. Fără Facade, fiecare',
        'apelant ar trebui să construiască şi să',
        'ordoneze manual lanțul.',
    ]
)

# 7. ADAPTER
pattern_slide(
    7, 'Adapter', 'Structural',
    'Adapter_FfmpegAdapter.png',
    integration_text=[
        'FfmpegAdapter implementează IFfmpegTool',
        'prin lansarea FFmpeg ca subprocess:',
        '',
        'ffmpeg -ss 30 -t 30 -i <file>',
        '       -f f32le -ac 1 -ar 44100 <tmp>',
        '',
        'Returnează float[] (eşantioane PCM)',
        'către analizatorul spectral.',
    ],
    problem_text=[
        'Analizatoarele spectrale au nevoie de PCM brut.',
        'FFmpeg este un binar extern cu CLI proprie.',
        'Adapter-ul ascunde invocarea procesului',
        'şi conversia bytes → float[].',
    ]
)

# 8. PROXY
pattern_slide(
    8, 'Proxy', 'Structural',
    'Proxy_CachingUserRepository.png',
    integration_text=[
        'CachingUserRepositoryProxy intercalează',
        'IMemoryCache între UserService şi SQLite.',
        '',
        'Hit (TTL 30s) → returnează imediat.',
        'Miss → delegă la SqliteUserRepository,',
        '        stochează rezultatul în cache.',
        '',
        'UpdateAsync() invalidează intrarea —',
        'statusul de ban/timeout rămâne proaspăt.',
    ],
    problem_text=[
        'Fiecare mesaj Telegram declanşează o verificare',
        'de moderare. Fără cache → câte un round-trip',
        'SQLite per mesaj, chiar dacă statusul nu s-a',
        'schimbat în ultimele 30 de secunde.',
    ]
)

# 9. STATE
pattern_slide(
    9, 'State', 'Comportamental',
    'State_AlbumSession.png',
    integration_text=[
        'AlbumSession delegă toate interacțiunile',
        'la IAlbumState curentă:',
        '',
        'IdleState → procesează fişiere individual,',
        '  /album-upload → tranziție la AlbumUploadState.',
        '',
        'AlbumUploadState → acumulează în PendingFiles,',
        '  /album-done → validare batch + Builder report',
        '              → tranziție înapoi la IdleState.',
    ],
    problem_text=[
        'Comportamentul sesiunii diferă radical după',
        'starea curentă. Fără State, AlbumSession ar',
        'conține lanțuri de if/switch fragile şi',
        'greu de extins.',
    ]
)

# 10. CHAIN OF RESPONSIBILITY
pattern_slide(
    10, 'Chain of Responsibility', 'Comportamental',
    'ChainOfResponsibility_ValidationPipeline.png',
    integration_text=[
        'Lanț: Extension → MIME → Spectral',
        '',
        'ExtensionHandler: respinge extensii',
        '  necunoscute sau lossy (mp3, aac, ogg...).',
        '',
        'MimeHandler: validează tipul MIME dacă',
        '  este prezent în metadatele Telegram.',
        '',
        'SpectralHandler: FFT via FFmpeg + analiza',
        '  energiei în banda 19–22 kHz.',
    ],
    problem_text=[
        'Fiecare etapă de validare are responsabilitate',
        'unică şi poate opri lanțul fără a cunoaşte',
        'etapele următoare. Handler-ele noi se adaugă',
        'fără a modifica codul existent.',
    ]
)

# 11. COMMAND
pattern_slide(
    11, 'Command', 'Comportamental',
    'Command_BotCommands.png',
    integration_text=[
        'Fiecare acțiune (/start, /ban, /album-done...)',
        'este un obiect IBotCommand cu ExecuteAsync().',
        '',
        'CommandDispatcher (Invoker) apelează',
        'ExecuteAsync() fără a cunoaşte comanda.',
        '',
        'HelpCommand verifică IsAdmin() şi afişează',
        'secțiunea admin doar administratorilor.',
    ],
    problem_text=[
        'Dispatcher-ul trebuie să rute mesajele fără',
        'a şti ce face fiecare comandă. Command',
        'decuplează complet routing-ul de execuție',
        'şi permite decorarea transparentă.',
    ]
)

# 12. STRATEGY
pattern_slide(
    12, 'Strategy', 'Comportamental',
    'Strategy_HashStrategy.png',
    integration_text=[
        'LedgerService (Context) deține IHashStrategy',
        'injectată la pornire din configurație:',
        '  Ledger:HashStrategy = "MD5" | "SHA256"',
        '',
        'IsDuplicateAsync() şi RecordAsync() apelează',
        'hashStrategy.ComputeAsync() — indiferent',
        'de algoritmul ales.',
        '',
        'Schimbarea algoritmului = 0 linii de cod.',
    ],
    problem_text=[
        'Algoritmul de deduplicare trebuie să fie',
        'configurabil fără recompilare.',
        'Strategy permite înlocuirea MD5 cu SHA-256',
        '(sau orice alt algoritm viitor) transparent',
        'față de LedgerService.',
    ]
)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE — DEMONSTRAREA SISTEMULUI
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
accent_bar(s)

txbox(s, 'Demonstrarea sistemului', Inches(0.5), Inches(0.15), W - Inches(1), Inches(0.7),
      size=30, bold=True, color=C_LIGHT)

# Big video placeholder box
vx, vy = Inches(0.5), Inches(1.1)
vw, vh = W - Inches(1), Inches(5.3)
rect(s, vx, vy, vw, vh, RGBColor(0x0D, 0x14, 0x26))
rect(s, vx, vy, vw, Inches(0.04), C_ACCENT)
rect(s, vx, vy + vh - Inches(0.04), vw, Inches(0.04), C_ACCENT)

txbox(s, '▶', vx + vw/2 - Inches(0.6), vy + Inches(1.5), Inches(1.2), Inches(1.2),
      size=72, color=C_ACCENT, align=PP_ALIGN.CENTER)
txbox(s, 'Video demonstrație', vx, vy + Inches(2.9), vw, Inches(0.6),
      size=22, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)
txbox(s,
      'Inserați aici videoclipul demonstrativ al aplicației NurFlac',
      vx, vy + Inches(3.55), vw, Inches(0.5),
      size=14, color=RGBColor(0x60, 0x60, 0x70), align=PP_ALIGN.CENTER, italic=True)

# Bullet points below (what the demo will show)
demos = [
    'Upload fişier FLAC autentic → acceptat ✓',
    'Upload fişier MP3 recodificat în FLAC → respins (spectral failure) ✗',
    'Sesiune /album-upload cu mai multe fişiere → raport final',
    'Comenzi admin: /ban, /timeout, /unban',
]
dy = vy + vh + Inches(0.12)
for i, d in enumerate(demos):
    bx_off = Inches(0.5) + (i % 2) * (vw / 2 + Inches(0.1))
    by_off = dy + (i // 2) * Inches(0.28)

rect(s, Inches(0.5), dy - Inches(0.06), vw, Inches(0.68),
     RGBColor(0x10, 0x18, 0x30))

demo_shape = slide_shapes = s.shapes
txbox(s, '  ·  '.join(demos[:2]),
      Inches(0.65), dy, vw - Inches(0.3), Inches(0.28),
      size=12, color=C_MUTED)
txbox(s, '  ·  '.join(demos[2:]),
      Inches(0.65), dy + Inches(0.3), vw - Inches(0.3), Inches(0.28),
      size=12, color=C_MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE FINAL — CONCLUZII
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)

# Large background accent
rect(s, 0, 0, Inches(0.5), H, C_ACCENT)

accent_bar(s, y=Inches(0.08), thickness=Inches(0.055))

txbox(s, 'Concluzii', Inches(0.7), Inches(0.25), W - Inches(1), Inches(0.75),
      size=36, bold=True, color=C_LIGHT)

rect(s, Inches(0.7), Inches(1.05), Inches(8.5), Inches(0.04), C_ACCENT)

conclusions = [
    ('12 şabloane GoF', 'fiecare rezolvă o problemă concretă, nu ornamentală'),
    ('Analiză spectrală FFT', 'detectează recodificări lossy invizibile la ochi liberi'),
    ('Moderare bazată pe State', 'escaladare automată: 3 strike → 24h, 4 → 72h, 5+ → ban'),
    ('Decorator mutual exclusiv', 'comenzile admin bypass-ează moderarea — design deliberat'),
    ('Proxy + Cache', 'reduce round-trip-urile SQLite la 1 per 30s per utilizator'),
    ('Extensibilitate', 'noi formate, strategii hash sau handler-e = 0 modificări existente'),
]

cy = Inches(1.25)
for title, body in conclusions:
    rect(s, Inches(0.7), cy, Inches(11.7), Inches(0.62), C_CARD)
    rect(s, Inches(0.7), cy, Inches(0.08), Inches(0.62), C_ACCENT)
    txbox(s, title, Inches(0.9), cy + Inches(0.04), Inches(3.2), Inches(0.3),
          size=14, bold=True, color=C_ACCENT)
    txbox(s, body, Inches(4.2), cy + Inches(0.04), Inches(8.1), Inches(0.52),
          size=13, color=C_LIGHT)
    cy += Inches(0.72)

# Final message
rect(s, Inches(0.7), cy + Inches(0.1), Inches(11.7), Inches(0.65),
     RGBColor(0x12, 0x0A, 0x1A))
txbox(s,
      'Șabloanele de proiectare nu sunt teorie — sunt soluții dovedite la probleme reale.',
      Inches(0.9), cy + Inches(0.18), Inches(11.3), Inches(0.4),
      size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)

prs.save(OUT)
print(f'Presentation saved to: {OUT}')
